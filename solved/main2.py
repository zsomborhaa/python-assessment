import collections 
import collections.abc
import json
import pathlib
import csv
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
import numpy as np
import logging
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


slidescreated = 0 #to log the total number of slides

def cimes_dia(cim, tartalom, presentation):                        #using the python-pptx library, each method determines what type of slides shall be created
    uj_dia = presentation.slides.add_slide(presentation.slide_layouts[0]) 
    uj_dia.shapes.title.text = cim
    uj_dia.placeholders[1].text = tartalom


def szoveges_dia(cim, tartalom, presentation):
    uj_dia = presentation.slides.add_slide(presentation.slide_layouts[1])
    uj_dia.shapes.title.text = cim
    uj_dia.placeholders[1].text = tartalom


def kep_dia(title, picture_path, presentation):
    uj_dia = presentation.slides.add_slide(presentation.slide_layouts[5])
    uj_dia.shapes.title.text = title
    bal_pozicio = vertikalis_pozicio = Inches(2)
    uj_dia.shapes.add_picture(picture_path, bal_pozicio, vertikalis_pozicio)
    """
    A picture slide once again extracts information from the JSON file to set the title and the image is based on filename
    slide_layouts contains a list of slide types that potentially be created (not all of them allow image insertion)
    Arguments needed for this method:
        title - The main title for the slide
        picture_path - The path to the picture file. This file must exist and needs to be in the correct folder
    """

    





"""
    This plot slide function requires X and Y axis values in addition to a .dat file to utilise a Powerpoint feature to insert a plot graph based on data given 
    Series (or in simpler terms, a label) is added to the plot using chart_data.add_series to emphasize what the graph is representing 
    It was encouraged to use the Numpy library to help generate this plot, which will read the data and use a delimiter
    A line chart will be produced (XL_CHART_TYPE.LINE) that also sets the size and position of the graph and the data associated with it

       
    """

def grafikon_dia(title, data_file, x_label, y_label, presentation):
    data = np.genfromtxt(data_file, delimiter=';')   #the delimiter is set to the same character that separates the values in the .dat file, so x and y values can be assigned
    x_data = data[:, 0] #data extracted from column 0
    y_data = data[:, 1] #column 1

    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    chart_data = CategoryChartData()
    chart_data.categories = [str(val) for val in x_data]
    chart_data.add_series("ECON", [float(val) for val in y_data])

    left = Inches(2)  #an attempt to shift the graph in the middle of the slide
    top = Inches(1.5)

    width = Inches(7)  #adjusting the size of the graph
    height = Inches(5)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, left, top, width, height, chart_data
    )


def readoptions(fajlnev):
    with open(fajlnev) as file:
        configoptions = json.load(file)
    return configoptions


def lista_diaa(presentation, cim, tartalom):
        new_slide = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(new_slide)
        slide.shapes.title.text = cim
        for i, list_row in enumerate(tartalom):
            level = list_row["level"] - 1
            if i == 0:
                
                list_first_line = slide.placeholders[1]
                list_first_line.level = level
                list_first_line.text = list_row["text"]
            else:
                list_line = list_first_line.text_frame.add_paragraph()
                list_line.level = level
                list_line.text = list_row["text"]


presentation = Presentation()       #presentation object is assigned to the variable that can manage slides, text, shapes assoicated with the presentation
fajlnev = "presvalues.json"   #File should be in the same folder / directory as this class when executed, or the full directory should be provided e.g. C:/Users/....
configoptions = readoptions(fajlnev)


# Prompt user for slide creation based on available slide types. This allows any number of slides to be created in any order
while True:
    try:
        slidechoice = int(input("What type of PowerPoint slide do you want to create? (Enter 1: Title, 2: Text, 3: List, 4: Picture, 5: Plot, 0: Save File): "))
        if slidechoice not in range(0, 6):
            raise ValueError
        if slidechoice == 0:   #the while loop breaks when 0 is entered and the methods can no longer be called
          break
    except ValueError:
        print("The input should be a number from 0 to 5.")
        continue

    selectfromJson = configoptions["presentation"][slidechoice - 1]
    fajta = selectfromJson["type"]

    if fajta == "title":          #this method generates the text allocated to title and content 
        cimes_dia(selectfromJson["title"], selectfromJson["content"], presentation)
        slidescreated += 1
        
    elif fajta == "text":
        szoveges_dia(selectfromJson["title"], selectfromJson["content"], presentation)
        slidescreated += 1
        
    elif fajta == "list":
         lista_diaa(presentation, selectfromJson["title"], selectfromJson["content"])
         slidescreated += 1
         
         
    elif fajta == "picture":
         kep_dia(selectfromJson["title"], selectfromJson["content"], presentation)
         slidescreated += 1
         
    elif fajta == "plot":                 #within the json file,     if there is a type named exactly as "plot" (using the comparison operator)
        #then this method generates the plot based on the values in json file
        grafikon_dia(                            
            selectfromJson["title"],                    
            selectfromJson["content"],
            selectfromJson["configuration"]["x-label"],
            selectfromJson["configuration"]["y-label"],
            presentation
        )
        slidescreated += 1
    else:
        print("You have not selected a valid option")

    createdfile = "codeoutput.pptx"

try:
    presentation.save(createdfile)
    print(f"LOG: Presentation saved to {createdfile}.")        #A more proper implementation of logs would produce this kind of information in a text file
    print(f"LOG: Successfully created {slidescreated} slides.")  #other relevant examples for logs would be checking .dat or json file contains information that the program can work with

except Exception as e:                         #exception handling will instruct the user for a potential solution if the presentaton.save method does not work
    print(f"LOG: Error saving presentation: {e}")
    print(f"LOG: There might already be a file called {createdfile}. Close this presentation if it is open and restart this program")
    
